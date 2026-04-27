"""
Generate PowerPoint from kabipay-feature-matrix.csv.
Lives under kabipay-ui/project-documentation/ (client & stakeholder materials).
Requires: pip install python-pptx
"""
from __future__ import annotations

import csv
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
CSV_PATH = HERE / "kabipay-feature-matrix.csv"
OUT_PATH = HERE / "kabipay-feature-matrix.pptx"

ROWS_PER_SLIDE = 5


def load_rows() -> list[dict[str, str]]:
    with open(CSV_PATH, encoding="utf-8-sig", newline="") as f:
        reader = csv.DictReader(f)
        return [r for r in reader if any((v or "").strip() for v in r.values())]


def set_cell_text(cell, text: str, size_pt: float = 8, bold: bool = False) -> None:
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text or "—"
    r.font.size = Pt(size_pt)
    r.font.name = "Calibri"
    r.font.color.rgb = RGBColor(0x20, 0x20, 0x20)
    if bold:
        r.font.bold = True
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.0), Inches(1.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "KabiPay feature matrix"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.text = "UI status · database coverage · Keka (benchmark)"
    p2.font.size = Pt(20)
    p2.font.name = "Calibri"
    p2.space_before = Pt(8)
    p3 = tf.add_paragraph()
    p3.text = f"Data source: kabipay-ui FEATURES.md + kabipay-database (Liquibase). Folder: project-documentation. CSV: {CSV_PATH.name}"
    p3.font.size = Pt(12)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(0x50, 0x50, 0x50)
    p3.space_before = Pt(16)


def add_legend_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = "How to read the columns"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = "Calibri"

    bullets = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.2), Inches(5.5))
    tf = bullets.text_frame
    lines = [
        "KabiPay UI: implemented / partial / not in UI — from implementation-validated UI documentation.",
        "Database: Liquibase domain id or module (tenant schema) where the data model is provisioned per tenant.",
        "Keka benchmark: indicative product areas on keka.com (Core HR, payroll, attendance, hiring, performance, EX, analytics, workflow, PSA) — for positioning, not a line-by-line audit.",
        "Notes: product talking points, gaps, or roadmap phrasing for client meetings.",
    ]
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = line
        para.font.size = Pt(16)
        para.font.name = "Calibri"
        para.level = 0
        para.space_after = Pt(10)


def add_table_slide(
    prs: Presentation, chunk: list[dict[str, str]], part: int, total: int
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.55))
    tp = title.text_frame.paragraphs[0]
    tp.text = f"Feature matrix (slide {part} of {total})"
    tp.font.size = Pt(20)
    tp.font.bold = True
    tp.font.name = "Calibri"

    nrows = 1 + len(chunk)
    ncols = 6
    left = Inches(0.35)
    top = Inches(0.85)
    width = Inches(12.5)
    height = Inches(5.4)
    table = slide.shapes.add_table(nrows, ncols, left, top, width, height).table

    col_widths = [Inches(1.05), Inches(1.9), Inches(1.25), Inches(1.35), Inches(1.55), Inches(2.5)]
    for i, w in enumerate(col_widths):
        if i < len(table.columns):
            table.columns[i].width = w

    headers = [
        "Category",
        "Feature / capability",
        "KabiPay UI",
        "Database",
        "Keka (benchmark)",
        "Notes for client",
    ]
    for c, h in enumerate(headers):
        set_cell_text(table.cell(0, c), h, size_pt=9, bold=True)

    keys = [
        "Category",
        "Feature / capability",
        "KabiPay UI",
        "Database (kabipay-database)",
        "Keka benchmark (indicative)",
        "Notes for client",
    ]
    for r, row in enumerate(chunk, start=1):
        for c, k in enumerate(keys):
            v = (row.get(k) or "").strip()
            if len(v) > 280:
                v = v[:277] + "…"
            set_cell_text(table.cell(r, c), v, size_pt=7.5)


def main() -> int:
    if not CSV_PATH.is_file():
        print(f"Missing CSV: {CSV_PATH}", file=sys.stderr)
        return 1
    rows = load_rows()
    if not rows:
        print("CSV is empty", file=sys.stderr)
        return 1

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_legend_slide(prs)

    chunks: list[list[dict[str, str]]] = []
    for i in range(0, len(rows), ROWS_PER_SLIDE):
        chunks.append(rows[i : i + ROWS_PER_SLIDE])
    total = len(chunks)
    for i, ch in enumerate(chunks, start=1):
        add_table_slide(prs, ch, i, total)

    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH} ({1 + 1 + total} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
